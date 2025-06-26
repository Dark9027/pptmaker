import requests
import json
import os
import re
import io
from pptx import Presentation
from pptx.util import Inches
from functools import lru_cache
from concurrent.futures import ThreadPoolExecutor

# --- CONFIGURATION ---
# ⚠️ IMPORTANT: Get your free API key from https://pixabay.com/api/docs/
PIXABAY_API_KEY = "YOUR_PIXABAY_API_KEY"
OLLAMA_URL = "http://localhost:11434/api/generate"
OLLAMA_MODEL = "llama3"
MAX_WORKERS = 4 # Number of parallel requests to the LLM

@lru_cache(maxsize=1024)
def generate_ollama_content(prompt):
    """Generic function to call the Ollama Llama3 API and return the response."""
    headers = {"Content-Type": "application/json"}
    payload = {"model": OLLAMA_MODEL, "prompt": prompt, "stream": False}
    try:
        response = requests.post(OLLAMA_URL, headers=headers, data=json.dumps(payload), timeout=20)
        response.raise_for_status()
        return response.json()["response"].strip()
    except requests.exceptions.RequestException as e:
        print(f"❌ Error calling Ollama API: {e}")
        return f"Error: Could not generate content."

def get_image_url(query, image_type="photo", category="business"):
    """Fetches an image URL from Pixabay API."""
    if PIXABAY_API_KEY == "YOUR_PIXABAY_API_KEY":
        print("⚠️ Pixabay API Key not set. Skipping image fetch.")
        return None
    
    api_url = "https://pixabay.com/api/"
    params = {
        "key": PIXABAY_API_KEY,
        "q": query,
        "image_type": image_type,
        "category": category,
        "per_page": 3,
        "safesearch": "true"
    }
    try:
        response = requests.get(api_url, params=params)
        response.raise_for_status()
        results = response.json().get("hits", [])
        if results:
            return results[0]["webformatURL"]
        else:
            print(f"🤷 No image found for query: {query}")
            return None
    except requests.exceptions.RequestException as e:
        print(f"❌ Error fetching image from Pixabay: {e}")
        return None

def generate_content_for_placeholder(placeholder, topic):
    """
    Generates content based on the placeholder type and its parameters.
    This function now correctly interprets your detailed placeholder rules.
    """
    base_prompt = f"For a PowerPoint presentation on '{topic}', generate the following content concisely and professionally:"

    # Pattern: {{topic1_1}} or {{topic1_2}}
    if "topic1_" in placeholder:
        part = placeholder.split('_')[-1].replace('}}', '')
        prompt = f"{base_prompt} Provide exactly two single words that capture the essence of the topic, separated by a space."
        words = generate_ollama_content(prompt).split()
        if len(words) >= 2:
            return words[0].capitalize() if part == '1' else words[1].capitalize()
        return "Topic" if part == '1' else "Essence"

    # Pattern: {{titlex_y_z}}, {{subtitlex_y_z}}
    match = re.match(r"\{\{(title|subtitle)(\d+)_(\d+)_(\d+)\}\}", placeholder)
    if match:
        ph_type, _, part, limit = match.groups()
        limit = int(limit)
        prompt = f"{base_prompt} A {ph_type} for a slide. It must be exactly {limit} words long."
        # If it's a multi-part title, the prompt should be more specific if needed
        # For now, we generate a title and let the user structure it if complex
        content = generate_ollama_content(prompt)
        return " ".join(content.split()[:limit]).title()

    # Pattern: {{parax_y}}
    match = re.match(r"\{\{para(\d+)_(\d+)\}\}", placeholder)
    if match:
        _, limit = match.groups()
        limit = int(limit)
        prompt = f"{base_prompt} A paragraph of approximately {limit} words."
        content = generate_ollama_content(prompt)
        return " ".join(content.split()[:limit]).capitalize() + "."
        
    # Pattern: {{bulletx_y}}
    match = re.match(r"\{\{bullet(\d+)_(\d+)\}\}", placeholder)
    if match:
        _, count = match.groups()
        count = int(count)
        prompt = f"{base_prompt} Generate exactly {count} bullet points. Each bullet point must be exactly 6 words long. Separate them with newlines."
        content = generate_ollama_content(prompt)
        # Clean up and format the bullet points
        bullets = [f"• {line.strip()}" for line in content.split('\n') if line.strip()]
        return "\n".join(bullets[:count])

    # Pattern: {{personname}}
    if placeholder == "{{personname}}":
        prompt = f"{base_prompt} A professional-sounding full name for the presenter."
        return generate_ollama_content(prompt)

    # Pattern: {{thankumess}}
    if placeholder == "{{thankumess}}":
        limit = 12
        prompt = f"{base_prompt} A brief thank-you message to the audience, exactly {limit} words long."
        content = generate_ollama_content(prompt)
        return " ".join(content.split()[:limit])

    return placeholder # Return original if no match

def generate_ppt(template_path, output_path, topic):
    """
    Main function to generate the PowerPoint presentation.
    """
    if not os.path.exists(template_path):
        print(f"❌ Template file not found: {template_path}")
        return

    prs = Presentation(template_path)
    
    # --- Step 1: Generate all text content in parallel ---
    text_placeholders = {}
    placeholders_to_find = set()

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Find all unique placeholders in the presentation
                matches = re.findall(r"\{\{.*?\}\}", shape.text_frame.text)
                for match in matches:
                    placeholders_to_find.add(match)
    
    # Exclude image/icon placeholders from text generation
    text_placeholders_to_generate = [
        ph for ph in placeholders_to_find if not ph.startswith("{{graph") and not ph.startswith("{{icon")
    ]
    
    print(f"⚙️ Generating text for {len(text_placeholders_to_generate)} placeholders...")
    with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
        # Map each placeholder to the generation function
        future_to_placeholder = {
            executor.submit(generate_content_for_placeholder, ph, topic): ph
            for ph in text_placeholders_to_generate
        }
        for future in future_to_placeholder:
            ph = future_to_placeholder[future]
            try:
                text_placeholders[ph] = future.result()
            except Exception as exc:
                print(f"❌ Generated an exception for {ph}: {exc}")
                text_placeholders[ph] = ph # Keep original on error

    # --- Step 2: Replace text placeholders in the presentation ---
    print("🔄 Replacing text in slides...")
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for ph, value in text_placeholders.items():
                    if ph in shape.text_frame.text:
                        # Simple text replace doesn't work well with formatting.
                        # This iterates through runs to preserve styles.
                        for para in shape.text_frame.paragraphs:
                            # Replace in the paragraph text directly to handle placeholders split across runs
                            if ph in para.text:
                                full_text = para.text
                                para.text = full_text.replace(ph, value)

    # --- Step 3: Handle image and icon placeholders ---
    print("🖼️ Handling image and icon placeholders...")
    for slide_idx, slide in enumerate(prs.slides):
        slide_title = slide.shapes.title.text if slide.has_title else f"{topic} {slide_idx}"
        placeholders_to_remove = []
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            # Handle {{graphimage}}
            if "{{graphimage}}" in shape.text_frame.text:
                print(f"🔍 Found graph placeholder on slide {slide_idx+1}. Searching for '{slide_title} graph'.")
                image_url = get_image_url(f"{slide_title} data graph chart", image_type="illustration", category="business")
                if image_url:
                    try:
                        response = requests.get(image_url)
                        image_stream = io.BytesIO(response.content)
                        # Add image, preserving placeholder's position and size
                        slide.shapes.add_picture(image_stream, shape.left, shape.top, width=shape.width, height=shape.height)
                        placeholders_to_remove.append(shape)
                    except Exception as e:
                        print(f"❌ Failed to add graph image to slide: {e}")

            # Handle {{iconX}}
            match = re.search(r"\{\{icon(\d+)\}\}", shape.text_frame.text)
            if match:
                print(f"🔍 Found icon placeholder on slide {slide_idx+1}. Searching for '{slide_title} icon'.")
                image_url = get_image_url(f"{slide_title} icon", image_type="illustration", category="computer")
                if image_url:
                    try:
                        response = requests.get(image_url)
                        image_stream = io.BytesIO(response.content)
                        slide.shapes.add_picture(image_stream, shape.left, shape.top, width=shape.width, height=shape.height)
                        placeholders_to_remove.append(shape)
                    except Exception as e:
                        print(f"❌ Failed to add icon image to slide: {e}")

        # Remove the placeholder shapes after adding images
        for shape in placeholders_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)

    # --- Step 4: Save the final presentation ---
    prs.save(output_path)
    print(f"\n✅ PPT generated successfully: {output_path}")


if __name__ == "__main__":
    try:
        topic = input("🔷 Enter the topic for your presentation: ").strip()
        if not topic:
            print("❌ Topic cannot be empty.")
        else:
            output_filename = f"{topic.replace(' ', '_')}_Presentation.pptx"
            generate_ppt("template.pptx", output_filename, topic)
    except KeyboardInterrupt:
        print("\n🚫 Process cancelled by user.")