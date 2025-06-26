from pptx import Presentation
import re

# ✅ Mock dynamic content generator (can be replaced with LLaMA)
def generate_content(placeholder_type, topic, args=None):
    base_words = topic.lower().split() + ["intelligence", "systems", "automation", "future", "solutions", "insights", "tech", "performance", "management"]

    if placeholder_type == "topic":
        return base_words[:2]

    if placeholder_type == "title" or placeholder_type == "subtitle":
        part, limit = args
        words = base_words[part - 1:part - 1 + limit] if part > 0 else base_words[:limit]
        return " ".join(words[:limit]).title()

    if placeholder_type == "para":
        limit = args
        words = []
        while len(words) < limit:
            words.extend(base_words)
        return " ".join(words[:limit]).capitalize() + "."

    if placeholder_type == "bullet":
        count = args
        bullets = []
        for i in range(count):
            b = base_words[i:i + 6]
            while len(b) < 6:
                b += base_words
            bullets.append("- " + " ".join(b[:6]))
        return "\n".join(bullets)

    if placeholder_type == "thankumess":
        return f"Thank you for learning about {topic}! We hope it was insightful."

    if placeholder_type == "graphimage":
        return "[Graph Placeholder]"

    if placeholder_type == "icon":
        return "[Icon Placeholder]"

    if placeholder_type == "personname":
        return "Dyizan Intern"

    return f"{placeholder_type} for {topic}"

# ✅ Replace placeholders inside a shape
def replace_placeholders(slide, placeholder_map):
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text
                    for ph, value in placeholder_map.items():
                        if ph in text:
                            run.text = text.replace(ph, value.replace("\\n", "\n"))

# ✅ Main generator
def generate_ppt(template_path, output_path, topic):
    prs = Presentation(template_path)
    placeholder_map = {}

    topic_words = generate_content("topic", topic)
    placeholder_map["{{topic1_1}}"] = topic_words[0].capitalize()
    placeholder_map["{{topic1_2}}"] = topic_words[1].capitalize()

    patterns = {
        "title": r"\{\{title(\d+)_(\d+)_(\d+)\}\}",
        "subtitle": r"\{\{subtitle(\d+)_(\d+)_(\d+)\}\}",
        "para": r"\{\{para(\d+)_(\d+)\}\}",
        "bullet": r"\{\{bullet(\d+)_(\d+)\}\}",
        "thankumess": r"\{\{thankumess\}\}",
        "personname": r"\{\{personname\}\}",
        "graphimage": r"\{\{graphimage\}\}",
        "icon": r"\{\{icon(\d+)\}\}",
    }

    # Scan & generate content
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text = run.text

                    # Titles/Subtitles
                    for typ in ["title", "subtitle"]:
                        for match in re.finditer(patterns[typ], text, re.IGNORECASE):
                            x, y, z = map(int, match.groups())
                            ph = match.group(0)
                            content = generate_content(typ, topic, args=(y, z))
                            placeholder_map[ph] = content

                    # Paragraphs
                    for match in re.finditer(patterns["para"], text, re.IGNORECASE):
                        x, y = map(int, match.groups())
                        ph = match.group(0)
                        content = generate_content("para", topic, args=y)
                        placeholder_map[ph] = content

                    # Bullets
                    for match in re.finditer(patterns["bullet"], text, re.IGNORECASE):
                        x, y = map(int, match.groups())
                        ph = match.group(0)
                        content = generate_content("bullet", topic, args=y)
                        placeholder_map[ph] = content

                    # Static placeholders
                    for key in ["thankumess", "personname", "graphimage"]:
                        if re.search(patterns[key], text, re.IGNORECASE):
                            ph = f"{{{{{key}}}}}"
                            content = generate_content(key, topic)
                            placeholder_map[ph] = content

                    # Icons
                    for match in re.finditer(patterns["icon"], text, re.IGNORECASE):
                        ph = match.group(0)
                        placeholder_map[ph] = generate_content("icon", topic)

    # Replace in all slides
    for slide in prs.slides:
        replace_placeholders(slide, placeholder_map)

    prs.save(output_path)
    print(f"\n✅ PPT generated successfully: {output_path}")

# ✅ Example usage
if __name__ == "__main__":
    topic = input("🔷 Enter topic: ").strip()
    generate_ppt("template.pptx", f"{topic.replace(' ', '_')}_Presentation.pptx", topic)
